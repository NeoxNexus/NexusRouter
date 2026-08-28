#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
根据 docs/ppt/2026-08-27-结项评审PPT.md 的定稿文案生成结项评审 PPT。
用法: docs/ppt/.venv/bin/python docs/ppt/build_pptx.py
输出: docs/ppt/NexusRouter-结项评审.pptx
"""
import math
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "NexusRouter-结项评审.pptx")

# ---- 设计常量 ----
ACCENT = RGBColor(0x9E, 0x1B, 0x32)      # 酒红主色
ACCENT_DARK = RGBColor(0x7A, 0x12, 0x26)
DARK = RGBColor(0x33, 0x33, 0x33)        # 正文深灰黑
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BOX_FILL = RGBColor(0xF7, 0xF2, 0xF3)    # 淡酒红底
PLACEHOLDER = RGBColor(0xD9, 0xD9, 0xD9) # 截图占位灰
FONT = "Microsoft YaHei"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---- 基础工具 ----
def set_font(run, size=14, bold=False, color=DARK, italic=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    # 同步设置中文字体(east asian)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def add_text(slide, x, y, w, h, lines, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.1,
             space_after=4):
    """lines: list of str 或 list of (text, {size,bold,color}) 或 list of [runs...]"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        runs = line if isinstance(line, list) else [(line, {})]
        for text, kw in runs:
            r = p.add_run()
            r.text = text
            set_font(r, size=kw.get("size", size), bold=kw.get("bold", bold),
                     color=kw.get("color", color), italic=kw.get("italic", False))
    return tb


def add_box(slide, x, y, w, h, text, fill=BOX_FILL, line_color=ACCENT,
            size=13, bold=False, color=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            align=PP_ALIGN.CENTER, line_width=1.0):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_width)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.05
        runs = line if isinstance(line, list) else [(line, {})]
        for t, kw in runs:
            r = p.add_run()
            r.text = t
            set_font(r, size=kw.get("size", size), bold=kw.get("bold", bold),
                     color=kw.get("color", color))
    return sp


def add_arrow(slide, x, y, w, h, direction="right", color=ACCENT):
    shape = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
             "left": MSO_SHAPE.LEFT_ARROW, "up": MSO_SHAPE.UP_ARROW}[direction]
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_header(slide, title, page_no=None):
    # 左侧色块 + 标题 + 分隔线
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.34),
                                 Inches(0.09), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_text(slide, 0.78, 0.26, 12.0, 0.66, [title], size=25, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.06),
                                Inches(12.23), Pt(2.2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()
    ln.shadow.inherit = False
    if page_no is not None:
        add_text(slide, 12.35, 7.05, 0.75, 0.35, [f"{page_no} / 14"], size=10,
                 color=GRAY, align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def set_cell(cell, text, size=12, bold=False, color=DARK, fill=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = anchor
    cell.margin_left = cell.margin_right = Inches(0.07)
    cell.margin_top = cell.margin_bottom = Inches(0.03)
    tf = cell.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        lines = text.split("\n")
    elif all(isinstance(x, tuple) for x in text):
        lines = [text]  # 单行的多 run 列表
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.05
        runs = line if isinstance(line, list) else [(line, {})]
        for t, kw in runs:
            r = p.add_run()
            r.text = t
            set_font(r, size=kw.get("size", size), bold=kw.get("bold", bold),
                     color=kw.get("color", color))


def add_table(slide, x, y, w, h, data, col_widths, header_size=12, body_size=12,
              row_heights=None):
    """data: list of rows; row = list of cell(str 或 [runs])"""
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = gt.table
    for j, cw in enumerate(col_widths):
        table.columns[j].width = Inches(cw)
    if row_heights:
        for i, rh in enumerate(row_heights):
            table.rows[i].height = Inches(rh)
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            if i == 0:
                set_cell(table.cell(i, j), cell_text, size=header_size, bold=True,
                         color=WHITE, fill=ACCENT, align=PP_ALIGN.CENTER)
            else:
                fill = WHITE if i % 2 == 1 else RGBColor(0xF5, 0xEE, 0xEF)
                set_cell(table.cell(i, j), cell_text, size=body_size, fill=fill)
    return table


def placeholder(slide, x, y, w, h, label):
    sp = add_box(slide, x, y, w, h, [], fill=PLACEHOLDER,
                 line_color=RGBColor(0xA6, 0xA6, 0xA6),
                 shape=MSO_SHAPE.RECTANGLE, line_width=1.2)
    # 虚线边框
    ln = sp.line._get_or_add_ln()
    dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
    ln.append(dash)
    tf = sp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "🖼 " + label
    set_font(r, size=14, bold=True, color=RGBColor(0x59, 0x59, 0x59))
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "（灰色占位框，定稿前手动替换为真实截图）"
    set_font(r2, size=10, color=RGBColor(0x80, 0x80, 0x80))
    return sp


# =====================================================================
# P1 封面
# =====================================================================
s = prs.slides.add_slide(BLANK)
# 顶部与底部酒红色带
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background(); band.shadow.inherit = False
band2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18))
band2.fill.solid(); band2.fill.fore_color.rgb = ACCENT; band2.line.fill.background(); band2.shadow.inherit = False

add_text(s, 1.0, 2.35, 11.33, 1.3,
         ["NexusRouter —— 让每一次 AI 调用都落在合适的模型上"],
         size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.42), Inches(3.85), Inches(2.5), Pt(2.5))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background(); ln.shadow.inherit = False
add_text(s, 1.0, 4.1, 11.33, 0.6, ["面向中孚互联网大模型平台（new-api）的智能模型调度层"],
         size=20, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.35, 11.33, 0.8,
         ['"赛马争先"AI 创意大赛结项评审  |  端点产品研发中心',
          '汇报人：陈锡光  |  项目成员：陈锡光、徐宝酉、钟凯华  |  2026-08'],
         size=12, color=GRAY, align=PP_ALIGN.CENTER, space_after=2)
add_notes(s, """讲稿（约 20 秒）：
各位评委好。我汇报的项目是 NexusRouter，一个部署在中孚互联网大模型平台（new-api）前的智能模型调度层。它解决的问题，相信在座用过公司大模型平台的同事都有体感。""")

# =====================================================================
# P2 痛点
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, '痛点：模型很全，但缺一个"分诊台"', 2)
add_table(s, 0.55, 1.35, 12.23, 4.5, [
    ["现象", "现状", "后果"],
    [[("感冒也挂专家号", {"bold": True, "color": ACCENT, "size": 14})],
     "AI 编程 Agent 每分钟自动产生大量背景请求（读文件、确认状态、工具回调），复杂度天差地别，却 100% 打给同一个高档模型",
     "隐性浪费随使用规模持续放大"],
    [[("专家号全靠拼手速", {"bold": True, "color": ACCENT, "size": 14})],
     "GLM-5.2/5.1 有日/周限额，先到先得",
     "抢到的不一定是真正需要的任务，复杂任务可能被迫降级"],
    [[("一个科室排长队，\n隔壁门可罗雀", {"bold": True, "color": ACCENT, "size": 14})],
     "请求凭习惯集中在 DeepSeek-v4-flash / MiniMax3",
     "模型矩阵的价值没有充分释放"],
], col_widths=[2.6, 6.0, 3.63], header_size=14, body_size=13,
   row_heights=[0.55, 1.5, 1.1, 1.1])
add_box(s, 0.55, 6.25, 12.23, 0.7,
        [[("医院不缺科室，缺的是分诊台", {"bold": True, "color": ACCENT, "size": 15}),
          ("——模型矩阵不缺模型，缺的是调度员。", {"bold": True, "size": 15})]],
        fill=BOX_FILL, line_color=ACCENT)
add_notes(s, """讲稿（约 90 秒）：
先看现状。平台已经有了很完整的模型矩阵，从轻量的 MiniMax2.5 到旗舰 GLM-5.2，还有私有部署的 Qwen 系列——这就像一家科室齐全的医院。但现在的使用方式，相当于没有分诊台，病人进门直接自己找医生，于是出现三个现象：
第一，感冒也挂专家号。Claude Code、Cursor 这类 AI 编程助手是 Agent 形态，你给一个任务，它每分钟在后台自动产生大量请求：读文件、确认状态、工具回调。这些请求复杂度天差地别，却百分之百打给同一个高档模型——"继续"两个字和"重新设计这个模块的架构"，看的是同一个专家。这个浪费很隐蔽，单次几毛钱没人心疼，但乘以每个请求、每个人、每个月，就是可观的成本。
第二，专家号全靠拼手速。GLM-5.2 有日限额和周限额，谁能看上专家，取决于谁挂得早，而不是谁的病更重。
第三，一个科室排长队，隔壁门可罗雀。大家的请求凭使用习惯集中在个别模型上，模型矩阵的整体价值没有充分释放出来。
医院不缺科室，缺的是分诊台。我的项目，就是给公司的模型矩阵补上这个分诊台。""")

# =====================================================================
# P3 方案总览与接入
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "NexusRouter：部署在中孚互联网大模型平台前的智能调度层", 3)
# 左侧链路图
add_text(s, 0.55, 1.25, 4.6, 0.35, ["调用链路"], size=13, bold=True, color=ACCENT)
chain = [
    "AI 编程 Agent\nClaude Code / Cursor / Codex",
    "NexusRouter 智能调度层\n复杂度分类 · 按档路由",
    "中孚互联网大模型平台（new-api）\n配额 · 计费 · 模型网关",
    "模型矩阵\nM2.5 / M2.7 / DS-flash / GLM / 私有 Qwen",
]
cy = 1.7
for i, txt in enumerate(chain):
    fill = BOX_FILL if i != 1 else ACCENT
    color = DARK if i != 1 else WHITE
    add_box(s, 0.55, cy, 4.6, 1.05, txt, fill=fill,
            line_color=ACCENT if i != 1 else ACCENT_DARK,
            size=12, bold=(i == 1), color=color)
    if i < 3:
        add_arrow(s, 2.55, cy + 1.05, 0.6, 0.28, "down")
    cy += 1.33
# 右侧 三个零 + 两个自由
add_text(s, 5.6, 1.25, 7.2, 0.35, ['三个零 + 两个自由'], size=13, bold=True, color=ACCENT)
items = [
    ("零学习成本", "baseurl 指向路由、模型填 auto 即接入；个人配额、计费不变（API Key 透传）"),
    ("零接入改造", "OpenAI / Anthropic 双协议，Claude Code、Cursor、Codex 即插即用"),
    ("零安装门槛", "npm 一行命令安装（AI Agent 用户本就有 node/npm 环境）"),
    ("模型选择自由", "auto = 走智能路由；填具体模型名 = 直连该模型、旁路分类。路由是可选项不是强制绑定"),
    ("部署方式自由", "个人可客户端本地部署；公司场景推荐服务侧集中部署（nginx + NexusRouter + 平台）"),
]
iy = 1.7
for head, body in items:
    add_text(s, 5.6, iy, 7.2, 0.9,
             [[("● " + head + "：", {"bold": True, "color": ACCENT, "size": 14}),
               (body, {"size": 13})]], line_spacing=1.08)
    iy += 1.02
add_notes(s, """讲稿（约 90 秒）：
方案是在员工和平台之间加一个智能调度层。接入做到了"三个零"：零学习成本——baseurl 指向路由、模型填 auto 即可；零改造——双协议兼容主流 Agent；零门槛——npm 一行命令安装，用 AI 编程工具的同事本来就有 node 环境。
特别说明两个"自由"：第一，模型选择自由——填 auto 才走路由，填具体模型名就直连，路由只是一个可选项，不是强制绑定，用户随时可以切回去。第二，部署自由——个人可以本地跑，公司场景推荐服务侧集中部署，挂在平台前面统一运维。""")

# =====================================================================
# P4 先立标准：路由怎样才算"不影响工作"？
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, '先立标准：路由怎样才算"不影响工作"？', 4)
add_table(s, 0.55, 1.45, 12.23, 3.7, [
    ["标准", "含义", "为什么"],
    [[("① 响应时间可控", {"bold": True, "color": ACCENT, "size": 14})],
     "分类串行在转发之前，延迟 1:1 叠加到首字时间；路由前后体感必须零差异",
     "路由是每天高频使用的基础设施，慢一点点都会被立刻感知"],
    [[("② 复杂任务零欠档", {"bold": True, "color": ACCENT, "size": 14})],
     "判错代价不对称：欠档（复杂→低档）= 质量事故；过档（简单→高档）= 静默多花一点钱",
     "宁可过档，不可欠档"],
], col_widths=[2.3, 5.9, 4.03], header_size=14, body_size=13,
   row_heights=[0.55, 1.55, 1.55])
add_box(s, 0.55, 5.5, 12.23, 0.7,
        [[("这两条标准，决定了后面所有设计取舍。", {"bold": True, "size": 15})]],
        fill=BOX_FILL, line_color=ACCENT)
add_notes(s, """讲稿（约 60 秒）：
在讲方案之前，先立两条标准——这个项目做到什么程度算成功？
第一，响应时间可控。分类是串行在请求转发之前的，它的延迟会一比一叠加到首字时间上。路由是每天要承受成百上千次调用的基础设施，慢一点点都会被立刻感知，所以路由前后体感必须零差异。
第二，复杂任务零欠档。这里要注意判错的代价是不对称的：把复杂任务判简单了，是看得见的质量事故；把简单任务判复杂了，只是静默多花一点钱。所以标准不是抽象的"准确率"，而是一条底线——宁可过档，不可欠档。
接下来两页，就是这两条标准如何决定整个机制的设计。""")

# =====================================================================
# P5 核心机制：三级分类 × 四级调度（原 P4）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "毫秒级分类，让每个模型各得其所", 5)
add_text(s, 0.55, 1.22, 12.2, 0.35, ["三级分类机制"], size=13, bold=True, color=ACCENT)
add_table(s, 0.55, 1.58, 12.23, 2.2, [
    ["层级", "机制", "耗时", "说明"],
    [[("L0 规则层", {"bold": True, "color": ACCENT})],
     "中英双语高精度规则短路（问候/确认/工具回调等）", "~0.05ms", "大量背景请求直接识别"],
    [[("L1 启发式层", {"bold": True, "color": ACCENT})],
     "信号驱动的档位判定（代码块、推理词、任务规模等）", "<2ms", "升档必须由具体信号驱动"],
    [[("L2 语义层", {"bold": True, "color": ACCENT})],
     "调用公司私有部署的 Qwen3.6 做语义判官", "~百毫秒", "0 API 成本，疑难请求兜底"],
], col_widths=[1.7, 5.6, 1.5, 3.43], header_size=13, body_size=12.5,
   row_heights=[0.45, 0.58, 0.58, 0.58])
add_text(s, 0.55, 4.0, 12.2, 0.35, ["四级调度映射表"], size=13, bold=True, color=ACCENT)
add_table(s, 0.55, 4.36, 12.23, 2.55, [
    ["档位", "路由目标", "承接任务"],
    [[("SIMPLE", {"bold": True, "color": ACCENT})], "MiniMax2.5", "问候、确认、状态查询等背景请求"],
    [[("MEDIUM", {"bold": True, "color": ACCENT})], "MiniMax2.7", "常规编码、文件修改"],
    [[("COMPLEX", {"bold": True, "color": ACCENT})], "DeepSeek-v4-flash", "复杂分析、多步任务"],
    [[("REASONING", {"bold": True, "color": ACCENT})], "GLM-5.1/5.2",
     "架构设计、严格推理——旗舰限额只给真正需要的任务"],
], col_widths=[2.2, 3.0, 7.03], header_size=13, body_size=12.5,
   row_heights=[0.45, 0.52, 0.52, 0.52, 0.52])
add_notes(s, """讲稿（约 2 分钟）：
核心机制是"三级分类、四级调度"。这套级联首先是为第一条标准服务的：最快的规则层 0.05 毫秒先短路，能接住的请求根本走不到后面的层，延迟几乎为零。
分类是级联的：最快最便宜的规则层先短路——Agent 产生的大量背景请求，问候、确认、工具回调，0.05 毫秒内被识别，直接发往轻量模型。接不住的进入启发式层，根据代码块、推理关键词、任务规模等信号判定档位，整个过程不到 2 毫秒。真正拿不准的少数请求，才进入语义层，由大模型做判官。
这里有一个设计我特别想讲：语义层的判官模型，恰到好处地用上了公司私有部署的 Qwen3.6。分类这件事本身不消耗任何互联网模型的 API 额度——公司的私有模型由此成为整个调度体系的智能底座，全部模型资产被串联成一个有机整体。私有部署的边际成本为零，"舍不舍得调"的问题天然不存在。
分类的结果映射到四个档位：背景请求去 MiniMax2.5，常规编码去 MiniMax2.7，复杂分析去 DeepSeek-v4-flash，只有真正的架构设计和严格推理，才动用 GLM-5.2。这就把第一章说的"专家号拼手速"变成了"按需分配"——限额流向真正需要它的任务，而不是手快的人。""")

# =====================================================================
# P6 分类准确性（原 P5）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "分类准不准，靠什么保障？", 6)
add_text(s, 0.55, 1.22, 12.2, 0.35, ["设计原则三条"], size=13, bold=True, color=ACCENT)
principles = [
    ("信号驱动", "升档必须由具体信号驱动（代码块、推理词、任务规模），不允许“拿不准就升档”"),
    ("代价非对称", "错误偏向后悔小的方向——宁可多花钱，不让复杂任务降级"),
    ("中英双语", "中文请求与英文同权处理（CJK 词数折算、中文关键词双通路）"),
]
px = 0.55
for head, body in principles:
    add_box(s, px, 1.6, 3.95, 1.35,
            [[(head, {"bold": True, "size": 14, "color": ACCENT})], [(body, {"size": 11.5})]],
            fill=BOX_FILL, line_color=ACCENT)
    px += 4.14
add_text(s, 0.55, 3.15, 12.2, 0.35, ["自我反馈闭环"], size=13, bold=True, color=ACCENT)
loop = ["全量路由日志\n档位/层级/置信度/耗时", "定期分析\n分布与偏差", "发现系统性偏差",
        "修复 + 53 条中英\n回归测试门禁"]
lx, lw, lgap = 0.55, 2.72, 0.45
for i, txt in enumerate(loop):
    add_box(s, lx, 3.55, lw, 1.0, txt, fill=WHITE, line_color=ACCENT, size=11.5)
    if i < 3:
        add_arrow(s, lx + lw + 0.05, 3.9, lgap - 0.1, 0.28, "right")
    lx += lw + lgap
# 回环箭头
add_arrow(s, 11.6, 4.62, 0.5, 0.35, "down", color=GRAY)
add_text(s, 3.2, 4.62, 8.2, 0.35, ["↩ 回到日志验证，形成闭环"],
         size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_box(s, 0.55, 5.15, 12.23, 1.75,
        [[("实例 D-002：", {"bold": True, "color": ACCENT, "size": 13}),
          ("通过 165 条真实日志发现“无条件升档棘轮”缺陷——流量集体打到最高档、路由形同虚设 → 修复决策结构 → 建立 53 条回归测试门禁，此后每次改动自动校验。", {"size": 12.5})],
         [("每条分类规则背后都有日志证据或测试用例。", {"bold": True, "size": 13})]],
        fill=BOX_FILL, line_color=ACCENT, align=PP_ALIGN.LEFT)
add_notes(s, """讲稿（约 90 秒）：
分类准确性是这个路由的生命线，我们用两层机制保障。
第一层是设计原则：升档必须由具体信号驱动，不允许"拿不准就往高档打"；错误代价是非对称的——把简单任务判复杂只是多花钱，把复杂任务判简单是质量事故，所以拿不准时偏向后悔小的方向。
第二层是自我反馈闭环：每一次路由决策都落日志，我定期分析分布有没有异常。举个真实例子：上线初期我通过 165 条日志发现一个系统性缺陷——三处"拿不准就升档"的逻辑叠加，导致几乎所有流量都被打到最高档，路由器形同虚设。修复之后，我建立了 53 条中英文回归测试作为门禁，之后每一次改动都自动校验。所以这套分类逻辑里没有一条规则是拍脑袋定的，背后都有日志证据或测试用例。""")

# =====================================================================
# P7 落地与效果（原 P6，两张控制台截图占位）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, '已上线运行：流量从"扎堆"到"各归其位"', 7)
add_text(s, 0.55, 1.22, 6.0, 0.35, [[("使用前（7月）", {"bold": True, "size": 13, "color": ACCENT}),
                                      ("  消耗 100% 扎堆 DeepSeek-V4-Pro/Flash 与 MiniMax-M3", {"size": 11})]])
add_text(s, 6.95, 1.22, 6.0, 0.35, [[("使用后（8.20-8.26）", {"bold": True, "size": 13, "color": ACCENT}),
                                      ("  自然分布到 M2.5 / M2.7 / M3 / DS-flash 四档", {"size": 11})]])
placeholder(s, 0.55, 1.6, 5.95, 2.7, "此处贴：使用前控制台截图（7月）")
placeholder(s, 6.85, 1.6, 5.95, 2.7, "此处贴：使用后控制台截图（8月）")
add_text(s, 0.55, 4.42, 12.23, 0.6,
         [[("GLM 本周未触达——限额被抢光，恰恰印证：没有调度时，稀缺限额被先到先得的低价值请求消耗。",
            {"size": 11.5, "color": GRAY, "italic": True})]])
add_box(s, 0.55, 5.0, 12.23, 0.85,
        [[("量化收益（同期反事实口径）：", {"bold": True, "color": ACCENT, "size": 13.5}),
          ("同一批请求若全走原高档模型需 ¥1.14，实际 ¥0.58，", {"size": 13}),
          ("节省 49.4%", {"bold": True, "color": ACCENT, "size": 16}),
          ("（非跨周期对比，剔除工作量差异）", {"size": 13})]],
        fill=BOX_FILL, line_color=ACCENT, align=PP_ALIGN.LEFT)
add_text(s, 0.55, 6.05, 12.23, 0.9,
         [[("体验：", {"bold": True, "color": ACCENT, "size": 13}),
           ("日常开发无任何体感差异", {"size": 13})],
          [("交付物：", {"bold": True, "color": ACCENT, "size": 13}),
           ("✅ 本地部署运行中（实时大屏现场可演示）  ✅ 路由日志  ✅ 代码仓与提交记录", {"size": 13})]],
         space_after=3)
add_notes(s, """讲稿（约 100 秒）：
项目已真实上线运行。前后对比看模型分布：左边是 7 月，我的消耗百分之百扎堆在 DeepSeek 和 MiniMax-M3 上——这是全公司使用模式的缩影。右边是接入路由后本周的数据：请求自然分布到了四个档位，背景请求去 M2.5，常规编码去 M2.7，复杂任务去 DeepSeek-flash。
有评委可能注意到 GLM 本周没有出现——原因很现实：限额太抢手，没抢到。这恰恰反过来证明了这个项目的价值：连已经接入路由的我都抢不到 GLM，说明稀缺限额正在被先到先得的方式消耗。如果全公司统一走路由集中调度，GLM 就只会流向真正的推理任务。
量化收益方面，路由大屏用反事实口径实测：同一批请求按原高档模型重定价需 1.14 元，实际 0.58 元，节省 49.4%——同期同批请求对比，剔除了工作量差异。而体验上，模型分布铺开的同时，日常开发没有任何体感差异。
回到开头立的两条标准：响应时间——路由引入的延迟实测不到 1 毫秒，体感零差异；复杂任务零欠档——53 条回归测试锁定分类行为，上线以来无欠档事故。
路由服务目前本地部署运行，实时大屏、路由日志、代码仓全部可供评委查验。""")

# =====================================================================
# P8 线 A：效果可复现验证（四宫格 + 大屏截图占位）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "效果不是写出来的，是可以回溯验证的", 8)
cells = [
    ("路由日志", "每次决策记录档位/层级/置信度/耗时，全量落盘可查"),
    ("53 条回归测试门禁", "每次改动自动校验分类行为不劣化"),
    ("完整审计链", "代码仓 commit 历史 + 设计评审文档全程留痕"),
]
positions = [(0.55, 1.45), (0.55, 4.35), (6.85, 4.35)]
for (head, body), (cx, cy2) in zip(cells, positions):
    add_box(s, cx, cy2, 5.95, 2.6,
            [[(head, {"bold": True, "size": 16, "color": ACCENT})], [(body, {"size": 13})]],
            fill=BOX_FILL, line_color=ACCENT)
# 右上：大屏截图占位
placeholder(s, 6.85, 1.45, 5.95, 2.6, "此处贴：实时大屏截图（档位分布 / 成本 / 节省）")
add_notes(s, """讲稿（约 60 秒）：
评审要点里有一条"效果可否复现验证"，这一页专门回应。每一次路由决策都落日志，档位、置信度、耗时全量可查；实时大屏现场可演示（本地部署）；53 条中英文回归测试构成门禁，任何改动导致分类劣化都会立刻红灯；代码仓的 commit 历史和设计评审文档全程留痕。欢迎各位评委回溯审计。""")

# =====================================================================
# P9 线 B-1：开发方法论（原 P8）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, 'AI 写了 ~100% 的代码，但每一行都过了"门禁"', 9)
add_text(s, 0.55, 1.25, 12.23, 0.65,
         [[("本项目全程重度 Vibe Coding，并从中沉淀出开源工作流框架 ", {"size": 13.5}),
           ("NexusRhythm", {"bold": True, "color": ACCENT, "size": 14}),
           ("（github.com/NeoxNexus/NexusRhythm）", {"size": 12.5, "color": GRAY})]])
flow = ["想法", "Discovery 收敛", "SDD 设计", "红灯测试", "AI 实现",
        "三门禁\n类型/构建/测试", "独立评审", "记忆蒸馏"]
fx, fw, fgap = 0.55, 1.34, 0.19
for i, txt in enumerate(flow):
    hot = i in (5,)
    add_box(s, fx, 2.05, fw, 0.85, txt, size=10.5, bold=hot,
            fill=ACCENT if hot else BOX_FILL,
            color=WHITE if hot else DARK,
            line_color=ACCENT_DARK if hot else ACCENT)
    if i < 7:
        add_arrow(s, fx + fw + 0.01, 2.32, fgap - 0.02, 0.3, "right")
    fx += fw + fgap
add_text(s, 0.55, 3.25, 12.2, 0.35, ["三个关键机制"], size=13, bold=True, color=ACCENT)
mechs = [
    ("状态机驱动", "ROADMAP 管理阶段状态，禁止跳步"),
    ("多智能体分工", "architect 设计 / reviewer 评审 / debt-collector 清债，互相制衡"),
    ("记忆蒸馏闭环", "Journal → ADR → /distill → Rules，踩过的坑自动变成 AI 下次的规则"),
]
my = 3.65
for head, body in mechs:
    add_text(s, 0.55, my, 12.23, 0.9,
             [[("● " + head + "：", {"bold": True, "color": ACCENT, "size": 14}),
               (body, {"size": 13})]], line_spacing=1.08)
    my += 0.95
add_notes(s, """讲稿（约 90 秒）：
这个项目的代码接近百分之百由 AI 生成，但它不是"放养式"的 Vibe Coding。我在开发过程中沉淀了一套工作流框架 NexusRhythm，已开源。
核心是让 AI 在纪律内跑：每个阶段必须先写设计文档和失败测试，AI 才能动手实现；提交前强制过类型检查、构建、全量测试三道门禁；架构设计、代码评审、债务清理由不同的智能体角色分工，互相制衡；每次踩的坑会被蒸馏成规则文件，AI 下次会话自动加载——犯过的错不犯第二次。
人的角色从"写代码"变成了"定方向、审设计、守门禁"。""")

# =====================================================================
# P10 线 B-2：AI 产出与人的分工（原 P9）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "AI 负责产能，人负责判断", 10)
rows9 = [
    ("工具链", "Claude Code + NexusRhythm（自定义命令 / 多智能体 / 规则体系）"),
    ("AI 产出占比", "代码、测试、文档 ≈100% 由 AI 生成"),
    ("人的角色四件事", "定义需求、架构决策、评审设计、最终验收"),
    ("Dogfooding", "本项目后期就是用“它自己路由的请求”开发的——8 月的路由数据里就包含开发它自己的流量"),
]
ry = 1.45
for head, body in rows9:
    add_text(s, 0.55, ry, 12.23, 0.85,
             [[("● " + head + "：", {"bold": True, "color": ACCENT, "size": 14.5}),
               (body, {"size": 13.5})]], line_spacing=1.1)
    ry += 0.95
add_box(s, 0.55, 5.35, 12.23, 1.55,
        [[("实例：", {"bold": True, "color": ACCENT, "size": 13.5}),
          ("AI 早期曾产出 2546 行的“15 维评分器”，评审后判断其不符合线上链路实际需求，", {"size": 13}),
          ("未接入", {"bold": True, "size": 13}),
          ("，换成更朴素可控的三层级联方案——决策过程完整留痕于设计评审文档。", {"size": 13})]],
        fill=BOX_FILL, line_color=ACCENT, align=PP_ALIGN.LEFT)
add_notes(s, """讲稿（约 60 秒）：
量化一下 AI 的使用。工具链是 Claude Code 加 NexusRhythm 的自定义命令、多智能体和规则体系。AI 产出占比接近百分之百——代码、测试、文档全部由 AI 生成，人的角色收敛为四件事：定义需求、架构决策、评审设计、最终验收。
举个真实例子：项目早期 AI 曾一口气写出一个 2500 多行的"15 维评分器"，看上去很唬人，但评审后发现它不符合线上链路的实际需求，最终没有接入，换成了更朴素但可控的三层级联方案。这个决策完整留痕在设计评审文档里——AI 负责产能，人负责判断。
还有一个有意思的事实：这个项目后期就是用"它自己路由的请求"开发出来的——8 月的路由数据里，就包含着开发它自己的流量。
（注：Token 消耗量不在 PPT 体现——多厂商、外网开发，无法准确统计，不讲无法佐证的数字。）""")

# =====================================================================
# P11 创新点（原 P10）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "三个原创设计", 11)
add_table(s, 0.55, 1.45, 12.23, 5.2, [
    ["创新", "内容", "对比传统方案"],
    [[("中英双语\n分层级联分类", {"bold": True, "color": ACCENT})],
     "规则短路 → 启发式 → 语义判官，亚毫秒级、零 token 开销",
     "传统路由多依赖云端 embedding，慢两个数量级且分类本身收费"],
    [[("代价非对称的\n决策设计", {"bold": True, "color": ACCENT})],
     "升档必须由信号驱动；错误偏向“后悔小”的方向",
     "多数方案只有简单/复杂两级，无明确代价模型"],
    [[("反事实基线计价", {"bold": True, "color": ACCENT})],
     "“同样的请求若走原模型需花多少”实时重定价，节省金额可逐笔审计",
     "“省钱”通常停留在估算，无法回溯"],
], col_widths=[2.7, 5.03, 4.5], header_size=14, body_size=13,
   row_heights=[0.55, 1.55, 1.55, 1.55])
add_notes(s, """讲稿（约 60 秒）：
创新点讲三个。第一，中英双语的分层级联分类，亚毫秒级、零额外 token 开销——传统方案用云端 embedding 做分类，慢两个数量级，而且每次分类本身还要花钱。第二，代价非对称的决策设计：判错的代价是不对称的，简单判复杂只是多花钱，复杂判简单是质量事故，所以升档必须由具体信号驱动。第三，反事实基线计价——"节省"不是估算，是同一批请求按原模型逐笔重定价，每一分钱都可审计。""")

# =====================================================================
# P12 竞品横向对比（7 列，小字号，原 P11）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, '赛道已被商业验证，但"零成本+私有化"的位置是空的', 12)
add_table(s, 0.45, 1.35, 12.43, 4.9, [
    ["维度", "NexusRouter", "Not Diamond", "OpenRouter Auto", "RouteLLM (LMSYS)",
     "vLLM Semantic Router", "LiteLLM/new-api 等网关"],
    ["路由机制", [("本地级联分类", {"bold": True})], "云端训练模型", "云端分类器",
     "本地训练分类器", "BERT 分类器", "不做复杂度分类"],
    ["分类成本", [("0 token", {"bold": True})], "$0.05/百万路由 token", "SaaS 内含",
     "需调 OpenAI embedding", "需 GPU/CPU 推理", "—"],
    ["判定延迟", [("<1ms", {"bold": True})], "官方自曝 100-150ms", "未公开", "未公开",
     "毫秒级", "—"],
    ["私有化部署", [("npm 单进程，内网可跑", {"bold": True})], "仅企业版", "❌ 纯 SaaS",
     "可，但 2024 年起停更", "需 K8s+Envoy+vLLM 整套栈", "可，但只管负载均衡"],
    ["路由粒度", "4 级复杂度阶梯", "全模型池", "~30 任务类型×成本档", "仅强/弱二分类",
     "多信号", "部署级"],
], col_widths=[1.35, 2.13, 1.95, 1.8, 1.75, 1.8, 1.65],
   header_size=10.5, body_size=10,
   row_heights=[0.6, 0.8, 0.8, 0.7, 1.1, 0.9])
add_box(s, 0.45, 6.45, 12.43, 0.62,
        [[("商业路由验证了赛道价值，但都是云端收费黑盒；开源方案要么停更、要么部署太重；网关类产品根本不看请求内容。",
           {"bold": True, "size": 12})]],
        fill=BOX_FILL, line_color=ACCENT)
add_notes(s, """讲稿（约 90 秒）：
这个赛道不是我凭空想象的，商业上已经验证过了——OpenRouter 的自动路由引擎最早就采购自 Not Diamond，而 Not Diamond 按路由 token 收费，每百万 token 5 美分，官方自曝路由本身要加 100 到 150 毫秒延迟。
但对我们的场景，它们都不合适：纯 SaaS 进不了内网，按 token 收费意味着"为了省钱先花钱"。开源侧，LMSYS 的 RouteLLM 从 2024 年就停更了；vLLM 社区的 Semantic Router 很活跃，但要部署 K8s、Envoy 加整套推理栈，太重了。而 LiteLLM、new-api 这类网关做的是负载均衡和计费，根本不看请求内容——这正是 NexusRouter 部署在 new-api 前面的合理性：一个管"该用哪档模型"，一个管"配额计费"，是互补不是竞争。
NexusRouter 占的位置是：分类零成本、亚毫秒、私有化、双协议、为 new-api 生态原生设计。""")

# =====================================================================
# P13 可复制性与推广路径（原 P12）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "双交付物，开箱即用", 13)
add_box(s, 0.55, 1.4, 6.0, 2.4,
        [[("交付物 1（产品）：NexusRouter", {"bold": True, "size": 15, "color": ACCENT})],
         [("npm 一行安装 / Docker 一键部署", {"size": 13})],
         [("完整 new-api 服务侧部署方案（含 nginx 配置）", {"size": 13})],
         [("产品化潜力：", {"bold": True, "color": ACCENT, "size": 12}),
          ("new-api 是行业广泛使用的开源网关，任何用 new-api / one-api 的团队都能直接复用", {"size": 12})]],
        fill=BOX_FILL, line_color=ACCENT, align=PP_ALIGN.LEFT)
add_box(s, 6.8, 1.4, 6.0, 2.4,
        [[("交付物 2（方法论）：NexusRhythm", {"bold": True, "size": 15, "color": ACCENT})],
         [("AI 工程工作流框架", {"size": 13})],
         [("clone 即用、可注入任何项目", {"size": 13})]],
        fill=BOX_FILL, line_color=ACCENT, align=PP_ALIGN.LEFT)
add_text(s, 0.55, 3.95, 12.2, 0.35, ["推广路径三步走"], size=13, bold=True, color=ACCENT)
steps12 = [
    ("✅ 个人试点", "已完成，数据已验证", True),
    ("✅ 7 人小团队部署测试", "已完成", True),
    ("🔲 公司级部署", "服务侧挂在平台前，全员 auto 接入，GLM 限额按需分配", False),
]
sx = 0.55
for head, body, done in steps12:
    add_box(s, sx, 4.4, 3.85, 1.9,
            [[(head, {"bold": True, "size": 14, "color": ACCENT if done else WHITE})],
             [(body, {"size": 12, "color": DARK if done else WHITE})]],
            fill=BOX_FILL if done else ACCENT,
            line_color=ACCENT if done else ACCENT_DARK)
    if sx < 8:
        add_arrow(s, sx + 3.87, 5.2, 0.3, 0.3, "right")
    sx += 4.19
add_notes(s, """讲稿（约 60 秒）：
可复制性方面，项目是双交付物。产品 npm 一行命令安装、Docker 一键部署，服务侧部署方案连 nginx 配置都是现成的；方法论 NexusRhythm 可注入任何项目。
推广路径上，前两步已经完成：我个人试点跑通了数据，7 人的小团队也已经部署测试。下一步是公司级服务侧部署——挂在平台前面，全员无感接入，旗舰限额按需分配。而且 new-api 是行业广泛使用的开源网关，这套方案不止适合我们公司，任何用 new-api 或 one-api 的团队都能直接复用——它具备作为产品向外推广的潜力。""")

# =====================================================================
# P14 总结与展望（原 P13）
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "让每一次 AI 调用都落在合适的模型上", 14)
add_box(s, 0.55, 1.45, 12.23, 1.15,
        [[("用 <1ms 的本地调度，换公司模型资产的整体提效——已验证节省 ",
           {"bold": True, "size": 17}),
          ("49.4%", {"bold": True, "size": 22, "color": ACCENT})]],
        fill=BOX_FILL, line_color=ACCENT)
add_text(s, 0.55, 2.95, 12.23, 1.0,
         [[("近期迭代：", {"bold": True, "color": ACCENT, "size": 14.5}),
           ("任务级路由记忆、可观测性增强（结构化决策日志 / metrics 端点）", {"size": 13.5})]],
         line_spacing=1.1)
add_text(s, 0.55, 3.85, 12.23, 1.6,
         [[("未来方向：", {"bold": True, "color": ACCENT, "size": 14.5}),
           ("结合公司安全主业——基于路由层的大模型审计与敏感请求阻断", {"size": 13.5})],
          [("（敏感代码定向私有模型、数据不出内网；路由层是天然的流量卡口）",
            {"size": 12.5, "color": GRAY})]],
         line_spacing=1.15)
add_text(s, 0.55, 5.9, 12.23, 1.0,
         ["谢谢各位评委，欢迎审计代码、验证数据。"],
         size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_notes(s, """讲稿（约 50 秒）：
总结一句话：用不到 1 毫秒的本地调度，换公司模型资产的整体提效，已验证的节省是 49.4%。
展望上，近期会做任务级路由记忆和可观测性增强。更远一步，我想点一个方向：路由层是公司 AI 流量的天然卡口，未来可以在这个位置生长出大模型审计、敏感请求识别与阻断能力——让敏感代码定向走私有模型、数据不出内网。这和公司的安全主业是同一个方向。
我的汇报到此，谢谢各位评委，欢迎审计代码、验证数据。""")

prs.save(OUT)
print(f"OK: {OUT}  ({len(prs.slides)} slides)")
