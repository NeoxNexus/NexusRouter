#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""自检生成的 pptx：页数、每页标题、shape 数、notes 非空、粗略溢出检查。"""
import math
import os
import unicodedata

from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
PATH = os.path.join(HERE, "NexusRouter-结项评审.pptx")

prs = Presentation(PATH)


def text_width_pt(text, size_pt):
    w = 0.0
    for ch in text:
        if unicodedata.east_asian_width(ch) in ("W", "F", "A"):
            w += size_pt
        else:
            w += size_pt * 0.55
    return w


def check_frame(tf, box_w_in, box_h_in, where, warnings):
    box_w_pt = box_w_in * 72
    box_h_pt = box_h_in * 72
    need_pt = 0.0
    for p in tf.paragraphs:
        runs = [(r.text, (r.font.size.pt if r.font.size else 14)) for r in p.runs]
        if not runs:
            need_pt += 10
            continue
        size = max(s for _, s in runs)
        total = sum(text_width_pt(t, s) for t, s in runs)
        lines = max(1, math.ceil(total / max(box_w_pt * 0.92, 1)))
        need_pt += lines * size * 1.25 + (p.space_after.pt if p.space_after else 0)
    if need_pt > box_h_pt * 1.18:
        warnings.append("  [溢出风险] %s: 估算需 %.0fpt, 框高 %.0fpt" % (where, need_pt, box_h_pt))


def walk_shapes(shapes, slide_idx, warnings, depth=0):
    for sh in shapes:
        if sh.shape_type == 6:  # group
            walk_shapes(sh.shapes, slide_idx, warnings, depth + 1)
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = sh.text_frame.text.strip().replace("\n", " ")[:24]
            check_frame(sh.text_frame, sh.width / 914400, sh.height / 914400,
                        "P%d '%s'" % (slide_idx, txt), warnings)
        if sh.has_table:
            tbl = sh.table
            for ri, row in enumerate(tbl.rows):
                rh = row.height / 914400
                for ci, cell in enumerate(row.cells):
                    cw = tbl.columns[ci].width / 914400
                    if cell.text.strip():
                        check_frame(cell.text_frame, cw, rh,
                                    "P%d 表格[%d,%d]" % (slide_idx, ri, ci), warnings)


def get_title(slide):
    # 取页面上方 y<1.1in 且字号最大的文本；封面页(顶部无标题)退化为全文最大字号文本
    def scan(top_limit):
        best, best_size = None, 0
        for sh in slide.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            if top_limit is not None and (sh.top is None or sh.top / 914400 >= top_limit):
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    sz = r.font.size.pt if r.font.size else 14
                    if sz > best_size:
                        best_size = sz
                        best = sh.text_frame.text.strip()
        return best
    return scan(1.1) or scan(None) or "(无标题)"


total_warnings = []
print("文件: %s" % PATH)
print("页数: %d\n" % len(prs.slides))
for idx, slide in enumerate(prs.slides, 1):
    n_shapes = len(slide.shapes)
    has_notes = slide.has_notes_slide and bool(slide.notes_slide.notes_text_frame.text.strip())
    notes_len = len(slide.notes_slide.notes_text_frame.text.strip()) if has_notes else 0
    warnings = []
    walk_shapes(slide.shapes, idx, warnings)
    print("P%-2d shapes=%-2d notes=%s(%d字)  标题: %s" %
          (idx, n_shapes, "✓" if has_notes else "✗", notes_len, get_title(slide)))
    for w in warnings:
        print(w)
    total_warnings.extend(warnings)

print("\n== 汇总 ==")
print("页数齐全: %s" % ("是" if len(prs.slides) == 14 else "否"))
print("溢出风险项: %d" % len(total_warnings))
